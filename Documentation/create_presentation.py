#!/usr/bin/env python3
"""Generate SmartStock+ MCO 1 Presentation PowerPoint."""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──────────────────────────────────────────────────────────────
PRIMARY      = RGBColor(0x1B, 0x5E, 0x20)   # Dark green
PRIMARY_LT   = RGBColor(0x2E, 0x7D, 0x32)   # Medium green
ACCENT       = RGBColor(0x43, 0xA0, 0x47)   # Light green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x21, 0x21, 0x21)
DARK_GRAY    = RGBColor(0x42, 0x42, 0x42)
MED_GRAY     = RGBColor(0x61, 0x61, 0x61)
LIGHT_GRAY   = RGBColor(0xE0, 0xE0, 0xE0)
BG_LIGHT     = RGBColor(0xF5, 0xF5, 0xF5)
CODE_BG      = RGBColor(0x26, 0x32, 0x38)   # Dark blue-gray for code
CODE_FG      = RGBColor(0xE0, 0xE0, 0xE0)
HIGHLIGHT    = RGBColor(0xFF, 0xC1, 0x07)   # Amber accent
LAYER_UI     = RGBColor(0x1B, 0x5E, 0x20)
LAYER_VM     = RGBColor(0x2E, 0x7D, 0x32)
LAYER_REPO   = RGBColor(0x43, 0xA0, 0x47)
LAYER_DATA   = RGBColor(0x66, 0xBB, 0x6A)
HILT_ACCENT  = RGBColor(0xFF, 0x6F, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──────────────────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p

def add_paragraph(tf, text, size=16, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                   space_before=Pt(4), space_after=Pt(2), font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p

def add_bullet(tf, text, size=16, color=BLACK, level=0, bold=False, font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p

def slide_header(slide, title, subtitle=None):
    """Add a green header bar with title at the top of a content slide."""
    # Header bar
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PRIMARY)
    # Title text
    txBox = add_textbox(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.5))
    set_text(txBox.text_frame, title, size=28, color=WHITE, bold=True)
    if subtitle:
        txBox2 = add_textbox(slide, Inches(0.6), Inches(0.6), Inches(11), Inches(0.4))
        set_text(txBox2.text_frame, subtitle, size=14, color=RGBColor(0xC8, 0xE6, 0xC9))
    # Slide number
    sn = add_textbox(slide, Inches(12.2), Inches(0.2), Inches(0.8), Inches(0.4))
    set_text(sn.text_frame, "", size=12, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.RIGHT)

def code_block(slide, left, top, width, height, code_text, font_size=11):
    """Add a dark code block with monospace text."""
    shape = add_shape(slide, left, top, width, height, CODE_BG)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Consolas"
        run.font.color.rgb = CODE_FG
    return shape


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═════════════════════════════════════════════════════════���════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, PRIMARY)

# Logo
try:
    logo_path = r"C:\Users\Admin\Downloads\SmartStock+\SmartStock\app\src\main\res\drawable\smartstock_logo.png"
    slide.shapes.add_picture(logo_path, Inches(5.6), Inches(0.6), Inches(2.1), Inches(2.1))
except:
    pass

# App name
txBox = add_textbox(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(1))
set_text(txBox.text_frame, "SmartStock+", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Tagline
txBox = add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5))
set_text(txBox.text_frame, "Offline-First Inventory & Asset Management System", size=20, color=RGBColor(0xC8, 0xE6, 0xC9), alignment=PP_ALIGN.CENTER)

# Divider line
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.5), Inches(4.3), Pt(2))
div.fill.solid()
div.fill.fore_color.rgb = HIGHLIGHT
div.line.fill.background()

# Team info
txBox = add_textbox(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(0.4))
set_text(txBox.text_frame, "Team 2  |  Mobile Programming 2  |  Midterm MCO Presentation", size=16, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

# Members - left column
members_left = [
    "Jhon Clarence B. Rulona", "Cris Gerard O. Carpon", "Jasper Keith P. Teorica",
    "Mark Jacob S. Allero", "John Dennis Y. Chan", "Leopoldo C. Villaflores III",
    "Dustin Aaron D. Tocayon"
]
members_right = [
    "Aleah A. Cantiga", "Angel Mae S. Dura", "Heizel D. Panugaling",
    "Justin Jamaica P. Julaton", "Roselyn B. Deguino",
    "Lou Ariane Mae B. Delos Reyes", "Rynel Jay V. Vallejos"
]

txBox = add_textbox(slide, Inches(2.2), Inches(5.35), Inches(4.5), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, m in enumerate(members_left):
    if i == 0:
        set_text(tf, m, size=11, color=RGBColor(0xE8, 0xF5, 0xE9), alignment=PP_ALIGN.LEFT)
    else:
        add_paragraph(tf, m, size=11, color=RGBColor(0xE8, 0xF5, 0xE9), space_before=Pt(1), space_after=Pt(1))

txBox = add_textbox(slide, Inches(6.8), Inches(5.35), Inches(4.5), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, m in enumerate(members_right):
    if i == 0:
        set_text(tf, m, size=11, color=RGBColor(0xE8, 0xF5, 0xE9), alignment=PP_ALIGN.LEFT)
    else:
        add_paragraph(tf, m, size=11, color=RGBColor(0xE8, 0xF5, 0xE9), space_before=Pt(1), space_after=Pt(1))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Project Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Project Overview", "What SmartStock+ does and who it serves")

# Three cards
card_data = [
    ("What It Does", "An Android inventory & asset management\napp that lets organizations track equipment,\ntools, and supplies on mobile devices\n\u2014 even without internet."),
    ("Purpose", "Replace error-prone paper logs and\nspreadsheets with a reliable, portable,\noffline-first mobile solution that\nensures data accessibility anytime."),
    ("Target Users", "School laboratories\nIT departments\nStudent organizations\nSmall offices & asset custodians"),
]

for i, (title, body) in enumerate(card_data):
    left = Inches(0.6 + i * 4.2)
    card = add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(5.2), WHITE, LIGHT_GRAY)
    card.shadow.inherit = False
    # Card title bar
    title_bar = add_rect(slide, left, Inches(1.5), Inches(3.8), Inches(0.7), PRIMARY_LT)
    txBox = add_textbox(slide, left + Inches(0.2), Inches(1.55), Inches(3.4), Inches(0.6))
    set_text(txBox.text_frame, title, size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Card body
    txBox = add_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(3.2), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, body, size=16, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Project Goal
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Project Goal", "Building a scalable, offline-first mobile foundation")

goals = [
    ("Scalable Architecture", "MVVM + Repository pattern separating UI, business logic, and data layers"),
    ("Offline-First Readiness", "Room database ensures full functionality without internet connectivity"),
    ("Hilt Dependency Injection", "Google's recommended DI framework for clean, automatic dependency management"),
    ("Role-Based Access Control", "Admin and Staff roles with granular permission levels (6 permissions each)"),
    ("Cloud Sync Ready", "Architecture prepared for Supabase integration in MCO 2"),
]

for i, (title, desc) in enumerate(goals):
    top = Inches(1.5 + i * 1.12)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), top + Inches(0.05), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_LT
    circle.line.fill.background()
    txn = add_textbox(slide, Inches(0.8), top + Inches(0.07), Inches(0.55), Inches(0.5))
    set_text(txn.text_frame, str(i+1), size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Text
    txBox = add_textbox(slide, Inches(1.6), top, Inches(10.5), Inches(0.35))
    set_text(txBox.text_frame, title, size=20, color=BLACK, bold=True)
    txBox2 = add_textbox(slide, Inches(1.6), top + Inches(0.35), Inches(10.5), Inches(0.4))
    set_text(txBox2.text_frame, desc, size=15, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Midterm Scope
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Midterm Scope", "MCO 1 deliverables vs MCO 2 planned features")

# Included card
inc_card = add_shape(slide, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.6), WHITE, LIGHT_GRAY)
inc_bar = add_rect(slide, Inches(0.5), Inches(1.4), Inches(7.5), Inches(0.6), PRIMARY)
txBox = add_textbox(slide, Inches(0.7), Inches(1.42), Inches(7), Inches(0.55))
set_text(txBox.text_frame, "MCO 1 - Included", size=20, color=WHITE, bold=True)

included = [
    "Room Database v7 \u2014 6 entities, full migration chain (v1\u2013v7)",
    "Hilt Dependency Injection (KSP-based approach)",
    "9 Jetpack Compose screens with Material Design 3",
    "MVVM + Repository data flow pattern",
    "Authentication with role-based access (Admin/Staff)",
    "Item usage tracking with destination location",
    "Return item with reason field",
    "Item image capture (camera + gallery)",
    "Barcode / QR code scanning",
]
txBox = add_textbox(slide, Inches(0.8), Inches(2.15), Inches(6.8), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(included):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = "\u2713  " + item
    run.font.size = Pt(14)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Segoe UI"

# Not included card
exc_card = add_shape(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(5.6), WHITE, LIGHT_GRAY)
exc_bar = add_rect(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.6), MED_GRAY)
txBox = add_textbox(slide, Inches(8.5), Inches(1.42), Inches(4), Inches(0.55))
set_text(txBox.text_frame, "MCO 2 - Planned", size=20, color=WHITE, bold=True)

excluded = [
    "Supabase cloud API integration",
    "Remote data synchronization",
    "Usage reporting with date filters",
    "WorkManager background sync",
    "Enhanced barcode flow",
]
txBox = add_textbox(slide, Inches(8.6), Inches(2.15), Inches(3.8), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(excluded):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(6)
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = "\u25CB  " + item
    run.font.size = Pt(14)
    run.font.color.rgb = MED_GRAY
    run.font.name = "Segoe UI"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: App Architecture Overview (Diagram)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "App Architecture Overview", "MVVM + Repository Pattern with Hilt Dependency Injection")

layer_w = Inches(8.5)
layer_h = Inches(1.1)
layer_left = Inches(0.8)
arrow_w = Inches(0.4)

layers = [
    (LAYER_UI,   "UI LAYER",        "Jetpack Compose + Material Design 3",
     "9 Screens  |  collectAsStateWithLifecycle()  |  Adaptive Layouts"),
    (LAYER_VM,   "VIEWMODEL LAYER", "@HiltViewModel + @Inject constructor",
     "InventoryViewModel  |  DashboardViewModel  |  Role Checks  |  StateFlow"),
    (LAYER_REPO, "REPOSITORY LAYER","InventoryRepository(dao, cloudSyncDataSource)",
     "Single Source of Truth  |  Cloud Sync Ready  |  NoOp Stub for MCO 1"),
    (LAYER_DATA, "DATA LAYER",      "Room Database v7 (Provided by Hilt @Module @Singleton)",
     "InventoryDao  |  6 Entities  |  35 Operations  |  6 Migrations  |  Flow Queries"),
]

for i, (color, title, sub, detail) in enumerate(layers):
    top = Inches(1.45 + i * 1.5)
    # Layer box
    shape = add_shape(slide, layer_left, top, layer_w, layer_h, color)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Pt(15)
    tf.margin_top = Pt(6)
    tf.word_wrap = True
    set_text(tf, title, size=18, color=WHITE, bold=True)
    add_paragraph(tf, sub, size=12, color=RGBColor(0xE8, 0xF5, 0xE9), space_before=Pt(1), space_after=Pt(0))
    add_paragraph(tf, detail, size=11, color=RGBColor(0xC8, 0xE6, 0xC9), space_before=Pt(1), space_after=Pt(0))

    # Arrow between layers
    if i < 3:
        arrow_top = top + layer_h
        # Down arrow label
        txBox = add_textbox(slide, Inches(4.3), arrow_top + Inches(0.02), Inches(2), Inches(0.35))
        labels = ["StateFlow \u2193  \u2191 Events", "suspend \u2193  \u2191 Flow", "suspend \u2193  \u2191 Flow"]
        set_text(txBox.text_frame, labels[i], size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Hilt sidebar
hilt_left = Inches(9.7)
hilt_box = add_shape(slide, hilt_left, Inches(1.45), Inches(3.1), Inches(5.6), RGBColor(0xFF, 0xF3, 0xE0), HILT_ACCENT)
tf = hilt_box.text_frame
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
tf.word_wrap = True
set_text(tf, "HILT DI", size=20, color=HILT_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "Wires all layers\nat compile time", size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
add_paragraph(tf, "", size=8, color=DARK_GRAY, space_before=Pt(10))
hilt_items = [
    "@HiltAndroidApp",
    "@HiltViewModel",
    "@Inject constructor",
    "@Module",
    "@Provides",
    "@Singleton",
    "@ApplicationContext",
]
for item in hilt_items:
    add_paragraph(tf, item, size=12, color=HILT_ACCENT, bold=False, alignment=PP_ALIGN.CENTER,
                  space_before=Pt(4), space_after=Pt(1), font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Why Use Hilt?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Why Use Hilt?", "Google's recommended Dependency Injection framework for Android")

benefits = [
    ("Compile-Time Safety", "Catches dependency errors at build time, not runtime \u2014 avoids crashes in production"),
    ("Annotation-Driven", "@HiltViewModel, @Inject, @Module replace hundreds of lines of manual factory code"),
    ("Scoped Lifecycles", "Dependencies tied to Android component lifecycles (Application, Activity, ViewModel)"),
    ("Reduced Boilerplate", "No manual ViewModelProvider.Factory classes needed \u2014 Hilt handles creation"),
    ("Scalability", "Easy to add cloud sync, WorkManager, and new ViewModels in MCO 2"),
    ("Google-Recommended", "Official DI solution for modern Android development"),
]

for i, (title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.5 + row * 1.85)
    card = add_shape(slide, left, top, Inches(6.1), Inches(1.55), WHITE, LIGHT_GRAY)
    txBox = add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.5), Inches(0.35))
    set_text(txBox.text_frame, title, size=18, color=PRIMARY, bold=True)
    txBox2 = add_textbox(slide, left + Inches(0.3), top + Inches(0.55), Inches(5.5), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    set_text(tf2, desc, size=14, color=DARK_GRAY)

# Note at bottom
txBox = add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
set_text(txBox.text_frame, "Note: Hilt Gradle plugin was incompatible with AGP 9.0.1 \u2014 resolved using KSP-only approach",
         size=12, color=HILT_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Hilt Module
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Hilt Module (Providing Dependencies)", "AppModule \u2014 @Module @InstallIn(SingletonComponent::class)")

code = '''@Module
@InstallIn(SingletonComponent::class)
object AppModule {

    @Provides @Singleton
    fun provideDatabase(@ApplicationContext context: Context): AppDatabase {
        return AppDatabase.getDatabase(context)
    }

    @Provides
    fun provideInventoryDao(database: AppDatabase): InventoryDao =
        database.inventoryDao()

    @Provides @Singleton
    fun provideCloudSyncDataSource(): CloudSyncDataSource =
        NoOpCloudSyncDataSource()

    @Provides @Singleton
    fun provideRepository(
        inventoryDao: InventoryDao,
        cloudSyncDataSource: CloudSyncDataSource
    ): InventoryRepository =
        InventoryRepository(inventoryDao, cloudSyncDataSource)

    @Provides @Singleton
    fun provideConnectivityObserver(
        @ApplicationContext context: Context
    ): ConnectivityObserver = NetworkConnectivityObserver(context)
}'''

code_block(slide, Inches(0.5), Inches(1.35), Inches(8.5), Inches(5.7), code, font_size=13)

# Explanation sidebar
expl = add_shape(slide, Inches(9.3), Inches(1.35), Inches(3.6), Inches(5.7), WHITE, LIGHT_GRAY)
tf = expl.text_frame
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
tf.word_wrap = True
set_text(tf, "How It Works", size=18, color=PRIMARY, bold=True)
explanations = [
    "@Module tells Hilt this class provides dependencies",
    "@InstallIn(SingletonComponent) = app-wide singletons",
    "Database via factory pattern with 6 migrations inside",
    "CloudSyncDataSource stubbed as NoOp for MCO 1",
    "Repository receives DAO + cloud sync source",
    "ConnectivityObserver monitors network state",
    "@Singleton ensures one instance per dependency",
    "Full chain resolved at compile time",
]
for exp in explanations:
    add_bullet(tf, "\u2022  " + exp, size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Data Model (Entity)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Data Model (Entity)", "6 Room entities powering the inventory system")

code = '''@Entity(
    tableName = "inventory_items",
    indices = [Index(value = ["assetCode"], unique = true)]
)
data class InventoryItem(
    @PrimaryKey(autoGenerate = true) val id: Int = 0,
    val assetCode: String,
    val name: String,
    val description: String?,
    val category: String,
    val quantity: Int,
    val inUseQuantity: Int = 0,
    val condition: String,
    val status: String,
    val location: String?,
    val imageUri: String? = null,
    val createdAt: Long,
    val lastUpdated: Long
)'''

code_block(slide, Inches(0.5), Inches(1.35), Inches(6.5), Inches(5.3), code, font_size=13)

# Entity cards on right
entities = [
    ("InventoryItem", "Core table \u2014 asset code, quantity, condition, status, image"),
    ("ItemHistory", "Audit trail \u2014 FK cascade delete to InventoryItem"),
    ("ItemUsageRecord", "Checkout tracking \u2014 destination, usedBy, returnReason"),
    ("LocalUser", "Admin/Staff accounts \u2014 unique email index"),
    ("CategoryEntity", "Dynamic categories \u2014 seeded defaults + custom"),
    ("AssetStatusEntity", "Available, In-Use, Damaged, Retired"),
]

for i, (name, desc) in enumerate(entities):
    top = Inches(1.4 + i * 0.9)
    card = add_shape(slide, Inches(7.3), top, Inches(5.5), Inches(0.75), WHITE, LIGHT_GRAY)
    txBox = add_textbox(slide, Inches(7.5), top + Inches(0.05), Inches(5), Inches(0.3))
    set_text(txBox.text_frame, name, size=14, color=PRIMARY, bold=True, font_name="Consolas")
    txBox2 = add_textbox(slide, Inches(7.5), top + Inches(0.35), Inches(5), Inches(0.35))
    set_text(txBox2.text_frame, desc, size=11, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: DAO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "DAO (Data Access Object)", "35 operations \u2014 all reads return Flow for reactive updates")

code = '''@Dao
interface InventoryDao {

    @Query("SELECT * FROM inventory_items ORDER BY lastUpdated DESC")
    fun getAllItems(): Flow<List<InventoryItem>>

    @Query("SELECT * FROM inventory_items WHERE name LIKE '%' || :q || '%'
            OR assetCode LIKE '%' || :q || '%'")
    fun searchItems(q: String): Flow<List<InventoryItem>>

    @Insert(onConflict = OnConflictStrategy.REPLACE)
    suspend fun insertItem(item: InventoryItem)

    @Update
    suspend fun updateItem(item: InventoryItem)

    @Delete
    suspend fun deleteItem(item: InventoryItem)

    @Query("SELECT * FROM item_usage_records WHERE itemId = :itemId
            AND status = 'Active'")
    fun getActiveUsageRecordsByItem(itemId: Int): Flow<List<ItemUsageRecord>>
}'''

code_block(slide, Inches(0.5), Inches(1.35), Inches(8.5), Inches(5.7), code, font_size=12)

# Key points sidebar
expl = add_shape(slide, Inches(9.3), Inches(1.35), Inches(3.6), Inches(5.7), WHITE, LIGHT_GRAY)
tf = expl.text_frame
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
tf.word_wrap = True
set_text(tf, "Key Points", size=18, color=PRIMARY, bold=True)
points = [
    "35 total operations for full CRUD",
    "All reads return Flow<List<T>> for reactive UI updates",
    "suspend functions for writes (coroutine-safe)",
    "Search by name AND asset code for barcode support",
    "Active usage records filtered by status = 'Active'",
    "History retrieval by item with FK relationship",
    "User queries for authentication",
]
for pt in points:
    add_bullet(tf, "\u2022  " + pt, size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Repository Pattern
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Repository Pattern", "Single source of truth between UI and data layers")

code = '''class InventoryRepository(
    private val inventoryDao: InventoryDao,
    private val cloudSyncDataSource: CloudSyncDataSource =
        NoOpCloudSyncDataSource()
) {
    val allItems: Flow<List<InventoryItem>> = inventoryDao.getAllItems()
    val categoryNames: Flow<List<String>> = inventoryDao.getAllCategoryNames()
    val cloudSyncStatus: StateFlow<CloudSyncStatus> =
        cloudSyncDataSource.status

    suspend fun insertItem(item: InventoryItem): Long {
        val itemId = inventoryDao.insertItem(item)
        cloudSyncDataSource.queueUpsertItem(item.copy(id = itemId.toInt()))
        return itemId
    }

    suspend fun updateItem(item: InventoryItem): Int {
        val updatedCount = inventoryDao.updateItem(item)
        if (updatedCount > 0) cloudSyncDataSource.queueUpsertItem(item)
        return updatedCount
    }

    fun getActiveUsageRecordsByItem(itemId: Int):
            Flow<List<ItemUsageRecord>> {
        return inventoryDao.getActiveUsageRecordsByItem(itemId)
    }
}'''

code_block(slide, Inches(0.5), Inches(1.35), Inches(8.0), Inches(4.8), code, font_size=13)

# Benefits
benefits_box = add_shape(slide, Inches(8.8), Inches(1.35), Inches(4.0), Inches(4.8), WHITE, LIGHT_GRAY)
tf = benefits_box.text_frame
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
tf.word_wrap = True
set_text(tf, "Why Repository?", size=18, color=PRIMARY, bold=True)
repo_points = [
    "Hilt's AppModule provides DAO + CloudSyncDataSource",
    "ViewModels never touch the DAO directly",
    "CloudSyncDataSource is NoOp stub \u2014 swap for Supabase in MCO 2",
    "Write ops queue changes to cloud sync automatically",
    "Handles inventory, usage records, history & users",
]
for pt in repo_points:
    add_bullet(tf, "\u2022  " + pt, size=12, color=DARK_GRAY)

# Future callout
future = add_shape(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.75), RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0x42, 0xA5, 0xF5))
txBox = add_textbox(slide, Inches(0.8), Inches(6.45), Inches(11.8), Inches(0.65))
set_text(txBox.text_frame, "MCO 2 Ready: Add Supabase remote data source alongside DAO \u2014 ViewModel layer stays untouched",
         size=14, color=RGBColor(0x1565, 0xC0, 0xFF)[:3] if False else RGBColor(0x15, 0x65, 0xC0), bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: ViewModel
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "ViewModel", "@HiltViewModel with reactive StateFlow state management")

code = '''@HiltViewModel
class InventoryViewModel @Inject constructor(
    private val repository: InventoryRepository
) : ViewModel() {

    val categoryNames: StateFlow<List<String>> = repository.categoryNames
        .map { it.ifEmpty { InventoryReferenceData.defaultCategories } }
        .stateIn(viewModelScope, SharingStarted.WhileSubscribed(5000),
                 InventoryReferenceData.defaultCategories)

    private val _currentUserRole = MutableStateFlow(UserRole.Staff)
    val currentUserRole: StateFlow<UserRole> = _currentUserRole.asStateFlow()

    fun useItem(item: InventoryItem, amount: Int, location: String) {
        if (!_currentUserRole.value.canAdjustUsage) {
            _operationError.value = "Your role cannot check out items"
            return
        }
        // Validate, update item, create UsageRecord, log history
    }

    fun returnItem(item: InventoryItem, amount: Int, reason: String) {
        // Close active usage records, restore quantities, log history
    }
}'''

code_block(slide, Inches(0.5), Inches(1.35), Inches(8.5), Inches(5.9), code, font_size=12)

# Sidebar
expl = add_shape(slide, Inches(9.3), Inches(1.35), Inches(3.6), Inches(5.9), WHITE, LIGHT_GRAY)
tf = expl.text_frame
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
tf.word_wrap = True
set_text(tf, "Key Concepts", size=18, color=PRIMARY, bold=True)
vm_points = [
    "@HiltViewModel \u2014 Hilt creates & manages lifecycle",
    "StateFlow only \u2014 no LiveData",
    "stateIn() with WhileSubscribed(5000ms) for efficiency",
    "combine() merges search + filter + sort into one stream",
    "Role permission checks in ViewModel layer",
    "useItem: checkout with destination tracking",
    "returnItem: close records with reason",
]
for pt in vm_points:
    add_bullet(tf, "\u2022  " + pt, size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: UI (Screenshots placeholder)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "UI (Compose)", "9 Jetpack Compose screens with Material Design 3")

# Create 9 placeholder boxes for screenshots in a 3x3 grid
screen_names = [
    "Splash Screen", "Auth / Login", "Dashboard",
    "Inventory List", "Add / Edit Item", "Item Detail",
    "History", "Camera Scan", "Profile & Settings",
]

for i, name in enumerate(screen_names):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.4 + row * 1.95)
    # Placeholder box
    placeholder = add_shape(slide, left, top, Inches(3.9), Inches(1.7), WHITE, ACCENT)
    tf = placeholder.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    # Vertical centering
    tf.paragraphs[0].space_before = Pt(20)
    run = p.add_run()
    run.text = "[Screenshot]"
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT
    run.font.name = "Segoe UI"
    # Label below
    add_paragraph(tf, name, size=13, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER, space_before=Pt(4))

# Note at bottom
txBox = add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.35))
set_text(txBox.text_frame, "Replace placeholders with actual app screenshots",
         size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Data Flow with Hilt
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Data Flow with Hilt", "Example: 'Use Item' action from tap to UI update")

steps = [
    ("1", "User Action",     "User taps 'Use' button,\nenters quantity & destination"),
    ("2", "ViewModel",       "Checks role permissions\nValidates quantity\n@HiltViewModel injected"),
    ("3", "Repository",      "Updates item quantities\nCreates ItemUsageRecord\nInserts ItemHistory\n@Inject injected"),
    ("4", "Room Database",   "Persists all changes locally\nEmits updated Flow"),
    ("5", "UI Updates",      "StateFlow triggers\nrecomposition automatically\ncollectAsStateWithLifecycle"),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.3 + i * 2.6)
    # Step box
    box = add_shape(slide, left, Inches(2.0), Inches(2.3), Inches(3.5), WHITE, PRIMARY_LT)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(1.6), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    txn = add_textbox(slide, left + Inches(0.8), Inches(1.62), Inches(0.6), Inches(0.55))
    set_text(txn.text_frame, num, size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    txBox = add_textbox(slide, left + Inches(0.15), Inches(2.2), Inches(2.0), Inches(0.4))
    set_text(txBox.text_frame, title, size=16, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    txBox2 = add_textbox(slide, left + Inches(0.15), Inches(2.7), Inches(2.0), Inches(2.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    set_text(tf, desc, size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        arrow_left = left + Inches(2.3)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(3.5), Inches(0.3), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# Key takeaway
txBox = add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "Key: Hilt wires ViewModel \u2192 Repository \u2192 DAO at compile time. The entire chain from user tap to UI update is reactive and automatic.",
         size=15, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Demonstration
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

txBox = add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1))
set_text(txBox.text_frame, "Live Demonstration", size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

txBox = add_textbox(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(0.5))
set_text(txBox.text_frame, "SmartStock+ in Action", size=20, color=RGBColor(0xC8, 0xE6, 0xC9), alignment=PP_ALIGN.CENTER)

demo_items = [
    "\u2460  App Launch \u2192 Splash \u2192 Login (bottom nav hidden)",
    "\u2461  Admin Login \u2192 Dashboard with summary cards",
    "\u2462  Add Item with image capture + custom category",
    "\u2463  Search, filter, view item detail with photo",
    "\u2464  Use Item \u2192 destination location tracking",
    "\u2465  Return Item \u2192 reason field",
    "\u2466  Condition Update (New / Good / Fair / Poor)",
    "\u2467  Staff Login \u2192 restricted permissions demo",
    "\u2468  Offline functionality \u2014 everything works without internet",
]

txBox = add_textbox(slide, Inches(3), Inches(3.6), Inches(7.3), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(demo_items):
    if i == 0:
        set_text(tf, item, size=16, color=RGBColor(0xE8, 0xF5, 0xE9))
    else:
        add_paragraph(tf, item, size=16, color=RGBColor(0xE8, 0xF5, 0xE9), space_before=Pt(5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Key Achievements
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Key Achievements", "MCO 1 deliverables verified by QA team")

achievements = [
    ("Room DB v7", "6 entities, 6 migrations,\n35 DAO operations"),
    ("Hilt DI", "KSP-based injection,\nAGP 9.0.1 compatible"),
    ("Offline-First", "Full functionality\nwithout internet"),
    ("Auth + Roles", "Admin/Staff login,\n6 granular permissions"),
    ("Usage Tracking", "Checkout with destination,\nreturn with reason"),
    ("Image Capture", "Camera + gallery,\nCoil display"),
    ("Barcode Scan", "GMS scanner,\ninventory lookup"),
    ("Adaptive UI", "Bottom nav + rail,\nWindowSizeClass"),
    ("Dark Mode", "MD3 dynamic theming,\nlight/dark toggle"),
    ("9 Screens", "Reactive StateFlow,\ncollectAsStateWithLifecycle"),
]

for i, (title, desc) in enumerate(achievements):
    col = i % 5
    row = i // 5
    left = Inches(0.3 + col * 2.56)
    top = Inches(1.5 + row * 2.8)
    card = add_shape(slide, left, top, Inches(2.36), Inches(2.4), WHITE, LIGHT_GRAY)
    # Title
    txBox = add_textbox(slide, left + Inches(0.15), top + Inches(0.25), Inches(2.06), Inches(0.5))
    set_text(txBox.text_frame, title, size=16, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    txBox2 = add_textbox(slide, left + Inches(0.15), top + Inches(0.8), Inches(2.06), Inches(1.3))
    tf = txBox2.text_frame
    tf.word_wrap = True
    set_text(tf, desc, size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Challenges Encountered
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Challenges Encountered", "Technical hurdles and how we resolved them")

challenges = [
    ("Hilt + AGP 9.0.1", "Hilt Gradle plugin dropped BaseExtension API",
     "KSP-only approach \u2014 extend Hilt_SmartStockApp() directly"),
    ("Room Migration Chain", "7 schema versions with data preservation",
     "Careful ALTER TABLE SQL + thorough migration testing"),
    ("Flow in Coroutines", "collect() inside launch suspends indefinitely",
     "Used .first() for one-shot queries within coroutines"),
    ("Reactive State Composition", "Merging search + filter + sort + role streams",
     "Kotlin combine operator for race-condition-free merging"),
    ("Camera Integration", "FileProvider + runtime permissions + content URIs",
     "Proper provider config, permission launcher, cache directory"),
]

for i, (title, problem, solution) in enumerate(challenges):
    top = Inches(1.4 + i * 1.18)
    card = add_shape(slide, Inches(0.5), top, Inches(12.3), Inches(1.0), WHITE, LIGHT_GRAY)
    # Challenge title
    txBox = add_textbox(slide, Inches(0.8), top + Inches(0.05), Inches(3.0), Inches(0.35))
    set_text(txBox.text_frame, title, size=16, color=PRIMARY, bold=True)
    # Problem
    txBox2 = add_textbox(slide, Inches(3.8), top + Inches(0.05), Inches(4.2), Inches(0.35))
    set_text(txBox2.text_frame, "Problem: " + problem, size=12, color=MED_GRAY)
    # Solution
    txBox3 = add_textbox(slide, Inches(3.8), top + Inches(0.45), Inches(8.5), Inches(0.4))
    set_text(txBox3.text_frame, "Solution: " + solution, size=13, color=DARK_GRAY, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Preparation for Final Term
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_header(slide, "Preparation for Final Term (MCO 2)", "Building on the MCO 1 foundation")

mco2_items = [
    ("Cloud Sync", "Supabase integration for real-time cross-device data synchronization.\nRepository layer already abstracts data \u2014 just add remote source."),
    ("Usage Tracking Screen", "Dedicated screen for Staff to view all checked-out items,\nlocations, and who checked them out. Monthly/annual reporting with date filters."),
    ("Background Sync", "WorkManager for reliable offline-to-cloud sync.\nPending changes sync automatically when connectivity restored."),
    ("Enhanced Scanning", "Improved barcode/QR flow with auto-match\nand batch scanning support."),
    ("UI Polish", "Two-pane empty state improvements,\nCompose animations, and remaining instructor recommendations."),
]

for i, (title, desc) in enumerate(mco2_items):
    top = Inches(1.4 + i * 1.15)
    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top + Inches(0.1), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_LT
    circle.line.fill.background()
    txn = add_textbox(slide, Inches(0.6), top + Inches(0.12), Inches(0.5), Inches(0.45))
    set_text(txn.text_frame, str(i+1), size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    txBox = add_textbox(slide, Inches(1.3), top, Inches(3.0), Inches(0.35))
    set_text(txBox.text_frame, title, size=18, color=PRIMARY, bold=True)
    # Description
    txBox2 = add_textbox(slide, Inches(1.3), top + Inches(0.35), Inches(11.3), Inches(0.7))
    tf = txBox2.text_frame
    tf.word_wrap = True
    set_text(tf, desc, size=13, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Conclusion
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

txBox = add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(0.8))
set_text(txBox.text_frame, "Conclusion", size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

summary_items = [
    "Fully functional offline-first inventory management system",
    "Clean MVVM architecture with Hilt Dependency Injection",
    "Room Database v7 \u2014 6 entities, 35 DAO operations, 6 migrations",
    "9 Jetpack Compose screens with Material Design 3",
    "Role-based authentication with granular permissions",
    "Usage tracking with destination locations and return reasons",
    "Image capture, barcode scanning, and adaptive layouts",
    "Architecture ready for MCO 2 cloud integration",
]

txBox = add_textbox(slide, Inches(2.5), Inches(2.2), Inches(8.3), Inches(3.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(summary_items):
    if i == 0:
        set_text(tf, "\u2713  " + item, size=17, color=RGBColor(0xE8, 0xF5, 0xE9))
    else:
        add_paragraph(tf, "\u2713  " + item, size=17, color=RGBColor(0xE8, 0xF5, 0xE9), space_before=Pt(6))

# Thank you
txBox = add_textbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6))
set_text(txBox.text_frame, "Thank you. We are open for questions.", size=22, color=HIGHLIGHT, bold=True, alignment=PP_ALIGN.CENTER)

# Team credit
txBox = add_textbox(slide, Inches(1), Inches(6.9), Inches(11.3), Inches(0.4))
set_text(txBox.text_frame, "Team 2  |  SmartStock+  |  Mobile Programming 2", size=14, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
output_path = r"C:\Users\Admin\Downloads\SmartStock+\SmartStock\Documentation\SmartStock_MCO1_Presentation_v2.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
