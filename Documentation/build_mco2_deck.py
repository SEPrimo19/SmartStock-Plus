"""SmartStock+ MCO 2 final presentation (.pptx) — comprehensive redesign.

Blue brand theme matched to the new app icon. Follows
MCO-2-Presentation-Format.docx (7 sections) + 2 screenshot slides.
Content reflects the current build: multi-tenant RLS, bidirectional
Supabase sync with natural-key reconciliation, Supabase Storage image
sync, biometric login, fully functional Profile, offline-first MVVM.

Run:  python Documentation/build_mco2_deck.py
Out:  SmartStock+_MCO2_Presentation.pptx in the project root.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Brand palette (matched to the new blue app icon) -------------------
DEEP   = RGBColor(0x0A, 0x23, 0x5B)   # title bg / header band
PRIMARY= RGBColor(0x1D, 0x4E, 0xD8)   # card headers, number chips
ACCENT = RGBColor(0x38, 0xBD, 0xF8)   # dividers, highlights
LT1    = RGBColor(0xDB, 0xEA, 0xFE)   # subtitle on dark
LT2    = RGBColor(0x93, 0xC5, 0xFD)   # meta on dark
LT3    = RGBColor(0xE8, 0xF1, 0xFE)   # member names on dark
PAGEBG = RGBColor(0xF4, 0xF7, 0xFB)   # content bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x1F, 0x2A, 0x37)
BODY2  = RGBColor(0x55, 0x63, 0x74)
HEAD   = RGBColor(0x0F, 0x17, 0x2A)
CARDLN = RGBColor(0xDD, 0xE6, 0xF2)
FONT   = "Segoe UI"

LOGO = r"C:\Users\Admin\Downloads\SmartStock+\SmartStock\app\src\main\res\drawable\smartstock_logo.png"
OUT  = r"C:\Users\Admin\Downloads\SmartStock+\SmartStock\SmartStock+_MCO2_Presentation.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
TOTAL = 9


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame


def run(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    r.font.name = FONT
    return r


def rrect(slide, shape, l, t, w, h, fill, line=None, line_w=1.25):
    s = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def header(slide, kicker, title, page):
    bg(slide, PAGEBG)
    rrect(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.18, DEEP)
    rrect(slide, MSO_SHAPE.RECTANGLE, 0, 1.18, 13.333, 0.055, ACCENT)
    k = _tb(slide, 0.62, 0.16, 10.0, 0.34)
    run(k.paragraphs[0], kicker.upper(), 12, ACCENT, bold=True)
    t = _tb(slide, 0.60, 0.46, 10.6, 0.62)
    run(t.paragraphs[0], title, 27, WHITE, bold=True)
    pg = _tb(slide, 11.7, 0.42, 1.2, 0.5)
    pg.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run(pg.paragraphs[0], f"{page:02d} / {TOTAL:02d}", 13, LT2, bold=True)
    # footer
    ft = _tb(slide, 0.62, 7.06, 9.0, 0.34)
    run(ft.paragraphs[0], "SmartStock+   ·   Mobile Programming 2 — MCO 2 Final Presentation",
        9, BODY2)
    fr = _tb(slide, 10.8, 7.06, 2.0, 0.34)
    fr.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run(fr.paragraphs[0], "Team 2", 9, BODY2, bold=True)


def card(slide, l, t, w, h, title, lines, body_size=14, title_size=16):
    rrect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, WHITE, line=CARDLN, line_w=1.0)
    rrect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, 0.60, PRIMARY)
    ht = _tb(slide, l + 0.20, t + 0.07, w - 0.40, 0.48)
    run(ht.paragraphs[0], title, title_size, WHITE, bold=True)
    bt = _tb(slide, l + 0.22, t + 0.76, w - 0.44, h - 0.92)
    bt.word_wrap = True
    first = True
    for ln, lvl in lines:
        p = bt.paragraphs[0] if first else bt.add_paragraph()
        first = False
        p.space_after = Pt(5)
        if lvl == 0:
            run(p, ln, body_size, BODY)
        else:
            run(p, "•  ", body_size, ACCENT, bold=True)
            run(p, ln, body_size, BODY2)


def numbered(slide, items, y0=1.55, row_h=1.02):
    y = y0
    for i, (title, desc) in enumerate(items, 1):
        c = rrect(slide, MSO_SHAPE.OVAL, 0.78, y + 0.02, 0.54, 0.54, PRIMARY)
        c.text_frame.word_wrap = False
        cp = c.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        run(cp, str(i), 17, WHITE, bold=True)
        tt = _tb(slide, 1.55, y - 0.06, 11.2, 0.38)
        run(tt.paragraphs[0], title, 18, HEAD, bold=True)
        dt = _tb(slide, 1.55, y + 0.30, 11.3, 0.62)
        dt.word_wrap = True
        run(dt.paragraphs[0], desc, 13, BODY2)
        y += row_h


def placeholder(slide, l, t, w, h, caption):
    s = rrect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h,
              RGBColor(0xEA, 0xF1, 0xFB), line=PRIMARY, line_w=1.5)
    s.text_frame.word_wrap = True
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, caption, 13, PRIMARY, bold=True)


def dbox(slide, l, t, w, h, title, sub):
    rrect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, WHITE, line=PRIMARY, line_w=1.25)
    tf = _tb(slide, l + 0.06, t + 0.10, w - 0.12, h - 0.16)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, title, 12, DEEP, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    run(p2, sub, 9, BODY2)


def arrow(slide, l, t, w, h, shape, label=None):
    rrect(slide, shape, l, t, w, h, ACCENT)
    if label:
        lt = _tb(slide, l - 0.45, t - 0.40, w + 0.9, 0.32)
        lt.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(lt.paragraphs[0], label, 8, BODY2, italic=True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# === SLIDE 1 — TITLE ====================================================
s = prs.slides.add_slide(BLANK)
bg(s, DEEP)
rrect(s, MSO_SHAPE.RECTANGLE, 0, 7.20, 13.333, 0.30, ACCENT)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.92), Inches(0.55),
                         Inches(1.5), Inches(1.5))
t = _tb(s, 1.0, 2.20, 11.3, 1.05)
t.paragraphs[0].alignment = PP_ALIGN.CENTER
run(t.paragraphs[0], "SmartStock+", 50, WHITE, bold=True)
st = _tb(s, 1.0, 3.18, 11.3, 0.5)
st.paragraphs[0].alignment = PP_ALIGN.CENTER
run(st.paragraphs[0],
    "Offline-First Inventory & Asset Management System", 19, LT1)
rrect(s, MSO_SHAPE.RECTANGLE, 4.92, 3.84, 3.50, 0.035, ACCENT)
mt = _tb(s, 1.0, 3.96, 11.3, 0.38)
mt.paragraphs[0].alignment = PP_ALIGN.CENTER
run(mt.paragraphs[0],
    "Team 2   |   Mobile Programming 2   |   MCO 2 Final Presentation",
    15, LT2)

COL1 = ["Jhon Clarence B. Rulona", "Cris Gerard O. Carpon",
        "Jasper Keith P. Teorica", "Mark Jacob S. Allero",
        "John Dennis Y. Chan", "Leopoldo C. Villaflores III",
        "Dustin Aaron D. Tocayon"]
COL2 = ["Angel Mae S. Dura", "Heizel D. Panugaling",
        "Justin Jamaica P. Julaton", "Roselyn B. Deguino",
        "Lou Ariane Mae B. Delos Reyes", "Rynel Jay V. Vallejos"]
for x, col in ((2.55, COL1), (7.05, COL2)):
    tf = _tb(s, x, 4.46, 4.30, 1.95)
    for i, name in enumerate(col):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run(p, name, 11, LT3)
inf = _tb(s, 1.0, 6.55, 11.3, 0.4)
inf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(inf.paragraphs[0],
    "Instructor:  Mr. Aaron Jude Pael      ·      May 19, 2026",
    13, LT2, bold=True)
notes(s, "Open with the app running on the device. Greet Mr. Pael, "
         "introduce Team 2 and the course/section. One line: SmartStock+ "
         "is an offline-first Android inventory app with secure multi-tenant "
         "Supabase cloud sync.")

# === SLIDE 2 — PROJECT OVERVIEW =========================================
s = prs.slides.add_slide(BLANK)
header(s, "01 · Introduction", "Project Overview", 2)
card(s, 0.60, 1.45, 3.95, 5.35, "Purpose & Problem", [
    ("Paper logs and always-online tools fail where connectivity is "
     "unreliable — data is lost or duplicated.", 0),
    ("SmartStock+ replaces them with a reliable, offline-first app "
     "that syncs to the cloud automatically.", 0),
])
card(s, 4.69, 1.45, 3.95, 5.35, "Target Users", [
    ("Small businesses", 1),
    ("School & IT laboratories", 1),
    ("Field & maintenance teams", 1),
    ("Asset custodians of shared equipment", 1),
])
card(s, 8.78, 1.45, 3.95, 5.35, "Core Modules", [
    ("Inventory CRUD + asset codes", 1),
    ("Barcode / QR scan & checkout", 1),
    ("Usage reports (CSV / PDF)", 1),
    ("Dashboard charts & low-stock alerts", 1),
    ("Stock-take / cycle count", 1),
    ("Multi-tenant cloud sync & auth", 1),
])
notes(s, "Stress the offline-first value proposition: never blocked by a "
         "dead network — Room is the source of truth, the cloud is a sync "
         "layer that reconciles automatically.")

# === SLIDE 3 — ARCHITECTURE =============================================
s = prs.slides.add_slide(BLANK)
header(s, "02 · How it's built", "Architecture — Offline-First MVVM", 3)

dbox(s, 0.55, 1.62, 2.70, 1.22, "Jetpack Compose UI", "screens & state")
dbox(s, 3.78, 1.62, 2.70, 1.22, "ViewModel", "StateFlow / events")
dbox(s, 7.01, 1.62, 2.70, 1.22, "Repository", "single mediator")
dbox(s, 10.24, 1.62, 2.55, 1.22, "Room DB", "local source of truth")
arrow(s, 3.27, 2.06, 0.48, 0.32, MSO_SHAPE.RIGHT_ARROW)
arrow(s, 6.50, 2.06, 0.48, 0.32, MSO_SHAPE.RIGHT_ARROW)
arrow(s, 9.73, 2.06, 0.48, 0.32, MSO_SHAPE.RIGHT_ARROW, "instant write")

arrow(s, 8.20, 2.86, 0.30, 0.50, MSO_SHAPE.DOWN_ARROW)
dbox(s, 6.45, 3.42, 3.05, 1.18, "SyncWorker", "WorkManager · bidirectional")
dbox(s, 10.05, 3.42, 2.74, 1.18, "Supabase", "PostgREST · Auth · Storage · RLS")
arrow(s, 9.52, 3.82, 0.55, 0.32, MSO_SHAPE.RIGHT_ARROW, "HTTPS")

note = rrect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 5.00, 8.90, 1.78,
             RGBColor(0xEA, 0xF1, 0xFB))
nt = _tb(s, 0.78, 5.16, 8.42, 1.46); nt.word_wrap = True
run(nt.paragraphs[0], "Offline-first: ", 13, DEEP, bold=True)
run(nt.paragraphs[0],
    "writes hit Room instantly so the UI never blocks. SyncWorker "
    "reconciles with Supabase when online — last-write-wins on "
    "updatedAt, soft-deletes, and natural-key (asset_code / name) "
    "matching so a re-installed or second device never duplicates rows.",
    12, BODY)
p2 = nt.add_paragraph(); p2.space_before = Pt(6)
run(p2, "Hilt ", 12, DEEP, bold=True)
run(p2, "wires every layer; ViewModels never touch DAOs or the network.",
    12, BODY)
dbox(s, 9.62, 5.00, 3.17, 1.78, "Clean separation",
     "data · domain · ui layers")
notes(s, "Walk left→right then down. The UI is fully decoupled from the "
         "data source — swapping the No-Op for the Supabase adapter took "
         "zero UI changes. Mention natural-key reconciliation as the fix "
         "that keeps sync duplicate-free across devices.")

# === SLIDE 4 — CLOUD / API INTEGRATION ==================================
s = prs.slides.add_slide(BLANK)
header(s, "03 · Supabase backend", "Cloud / API Integration", 4)
card(s, 0.60, 1.45, 7.45, 5.35, "Supabase (REST API backend)", [
    ("PostgreSQL + PostgREST REST API + GoTrue auth + Storage + Edge "
     "Functions, via supabase-kt over Ktor", 0),
    ("Auth — email/password, session persistence, password reset, "
     "biometric unlock; Admin provisions Staff via an Edge Function", 0),
    ("Bidirectional sync — 6 tables, last-write-wins, soft-delete, "
     "connectivity-gated; natural-key reconciliation prevents duplicates", 0),
    ("Supabase Storage — item photos sync across devices, per-team "
     "scoped objects", 0),
    ("Multi-tenant isolation — per-team Row-Level Security via a "
     "my_team() SECURITY DEFINER function; each account sees only its "
     "own team's data", 0),
])
placeholder(s, 8.25, 1.45, 4.48, 5.35,
            "[ SCREENSHOT ]\n\nSupabase table editor\n(inventory_items rows)\n"
            "\n+ optional Storage\nbucket / sample JSON")
notes(s, "Supabase IS a REST API backend (PostgREST) + managed auth + "
         "object storage — it fully satisfies the 'Firebase/REST API' "
         "requirement. Show a Supabase table screenshot. Highlight the "
         "per-team RLS isolation as the security story.")

# === SLIDE 5 — UI DEMO ==================================================
s = prs.slides.add_slide(BLANK)
header(s, "04 · Live walkthrough", "UI Demo", 5)
card(s, 0.60, 1.45, 12.13, 5.35, "Demo Flow (live on device)", [
    ("Registration / Login — redesigned, scrollable, biometric option", 1),
    ("Dashboard — KPIs, low-stock chart, connectivity & sync badge", 1),
    ("CRUD — add an item with photo + barcode, edit, delete with Undo", 1),
    ("Scan — link a barcode, scan-to-checkout a unit", 1),
    ("Cloud sync — row + photo appear in Supabase; pull-to-refresh", 1),
    ("Profile — theme, security settings, CSV export, team management", 1),
    ("Multi-tenant — a fresh account sees an empty, isolated inventory", 1),
], body_size=16)
notes(s, "DEMO SCRIPT (under 4 min):\n"
         "1. Login -> register a fresh account (proves isolation: empty "
         "inventory, only itself in the team list).\n"
         "2. Add an item with a camera photo + scanned barcode -> save.\n"
         "3. Supabase dashboard -> refresh inventory_items -> row present "
         "with team_id; show the Storage bucket has the photo.\n"
         "4. Edit quantity below threshold -> low-stock notification.\n"
         "5. Delete an item -> tap UNDO.\n"
         "6. Profile -> Export Data (CSV share sheet) + Security Settings.\n"
         "Keep screenshots open as a fallback.")

# --- SLIDES 6 & 7 — DEMO SCREENSHOTS ------------------------------------
for pageno, tag, labs in (
    (6, "Screenshots (1 of 2)",
     ["Login / Registration", "Dashboard", "Inventory list"]),
    (7, "Screenshots (2 of 2)",
     ["Add / Edit item (CRUD)", "Cloud sync / Supabase", "Profile & Settings"]),
):
    s = prs.slides.add_slide(BLANK)
    header(s, "04 · Live walkthrough", "UI Demo — " + tag, pageno)
    x = 0.95
    for lab in labs:
        placeholder(s, x, 1.55, 3.78, 5.15, "[ SCREENSHOT ]\n\n" + lab)
        x += 4.05
    notes(s, "Portrait phone captures — paste each over its box and delete "
             "the outline. Remove this slide if you demo fully live.")

# === SLIDE 8 — CHALLENGES & LEARNINGS ===================================
s = prs.slides.add_slide(BLANK)
header(s, "05 · What we learned", "Challenges & Learnings", 8)
numbered(s, [
    ("Room ↔ REST synchronization",
     "Reconciling epoch-millis vs ISO-8601 and last-write-wins conflict "
     "resolution across six tables."),
    ("Multi-tenant data isolation",
     "An over-broad is_admin() RLS policy + an un-wiped local cache leaked "
     "data across accounts — fixed with strict per-team RLS and a cache "
     "wipe on account switch."),
    ("Duplicate-free sync",
     "Client-UUID-only keying duplicated rows on re-install — fixed by "
     "reconciling items/categories/statuses by natural key before push."),
    ("Cross-device images",
     "Local file:// paths are meaningless on other devices — moved photos "
     "to per-team Supabase Storage with upload/download in the sync loop."),
    ("Best practices applied",
     "MVVM · Repository · StateFlow · Hilt DI · lifecycle-aware APIs · "
     "WorkManager · RLS-backed security · clean architecture."),
], y0=1.55, row_h=1.02)
notes(s, "The multi-tenant isolation bug is the strongest talking point — "
         "found in testing, root-caused to an is_admin() SELECT policy that "
         "gave every admin global read, fixed with strict per-team RLS.")

# === SLIDE 9 — CONCLUSION & FUTURE ======================================
s = prs.slides.add_slide(BLANK)
header(s, "06 · Wrap-up", "Conclusion & Future Improvements", 9)
card(s, 0.60, 1.45, 6.05, 4.95, "Summary", [
    ("A complete offline-first inventory system: full CRUD, scanning, "
     "reporting, charts, notifications, biometric auth, and secure "
     "multi-tenant Supabase sync with cross-device image storage.", 0),
    ("Stays fully usable with zero connectivity and reconciles "
     "automatically once back online — the core goal achieved.", 0),
])
card(s, 6.84, 1.45, 5.89, 4.95, "Future Improvements", [
    ("Hardware-backed biometrics — Keystore + BiometricPrompt "
     "CryptoObject encrypting the session token", 1),
    ("Admin web panel for fleet-wide oversight", 1),
    ("Real-time updates via Supabase Realtime", 1),
    ("FCM push notifications (currently device-local)", 1),
    ("Audit log + hardened team-bootstrap handling", 1),
])
ack = _tb(s, 0.60, 6.52, 12.13, 0.5)
ack.paragraphs[0].alignment = PP_ALIGN.CENTER
run(ack.paragraphs[0],
    "Thank you, Mr. Aaron Jude Pael, for your guidance throughout "
    "Mobile Programming 2.", 14, DEEP, bold=True, italic=True)
notes(s, "Close confidently: every MCO 2 requirement is met and the app is "
         "demo-ready. Read the acknowledgment aloud — thank Mr. Pael "
         "directly and invite his questions.")

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
